import csv
import io
import logging
from dataclasses import dataclass, field

logger = logging.getLogger(__name__)


@dataclass
class ParsedDocument:
    title: str
    sections: list[dict] = field(default_factory=list)


class DocumentParseException(Exception):
    pass


def _parse_txt(content: bytes, filename: str) -> ParsedDocument:
    text = content.decode("utf-8")
    return ParsedDocument(
        title=filename,
        sections=[{"level": 0, "title": filename, "content": text}],
    )


def _parse_pdf(content: bytes, filename: str) -> ParsedDocument:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(content))
    pages_text = []
    for i, page in enumerate(reader.pages):
        txt = page.extract_text() or ""
        pages_text.append(txt)
    full_text = "\n".join(pages_text)
    return ParsedDocument(
        title=filename,
        sections=[{"level": 0, "title": filename, "content": full_text}],
    )


def _parse_docx(content: bytes, filename: str) -> ParsedDocument:
    from docx import Document
    doc = Document(io.BytesIO(content))
    sections = []
    current_level = 0
    current_title = filename
    current_content = []
    for para in doc.paragraphs:
        style = para.style.name if para.style else ""
        if style.startswith("Heading"):
            if current_content:
                sections.append({
                    "level": current_level,
                    "title": current_title,
                    "content": "\n".join(current_content),
                })
            try:
                current_level = int(style.replace("Heading", "").strip())
            except ValueError:
                current_level = 0
            current_title = para.text
            current_content = []
        else:
            if para.text:
                current_content.append(para.text)
    if current_content or not sections:
        sections.append({
            "level": current_level,
            "title": current_title,
            "content": "\n".join(current_content),
        })
    return ParsedDocument(title=filename, sections=sections)


def _parse_xlsx(content: bytes, filename: str) -> ParsedDocument:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(content), data_only=True)
    sections = []
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        rows = []
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
            rows.append(row_text)
        sections.append({
            "level": 0,
            "title": sheet_name,
            "content": "\n".join(rows),
        })
    return ParsedDocument(title=filename, sections=sections)


def _parse_pptx(content: bytes, filename: str) -> ParsedDocument:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    sections = []
    for i, slide in enumerate(prs.slides):
        texts = []
        if slide.shapes.title and slide.shapes.title.text:
            texts.append(slide.shapes.title.text)
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text:
                        texts.append(paragraph.text)
        sections.append({
            "level": 0,
            "title": f"Slide {i + 1}",
            "content": "\n".join(texts),
        })
    return ParsedDocument(title=filename, sections=sections)


def _parse_md(content: bytes, filename: str) -> ParsedDocument:
    text = content.decode("utf-8")
    lines = text.split("\n")
    sections = []
    current_level = 0
    current_title = filename
    current_content = []
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("#"):
            if current_content or not sections:
                sections.append({
                    "level": current_level,
                    "title": current_title,
                    "content": "\n".join(current_content),
                })
            level = len(stripped) - len(stripped.lstrip("#"))
            current_level = level
            current_title = stripped.lstrip("#").strip()
            current_content = []
        else:
            current_content.append(line)
    if current_content or not sections:
        sections.append({
            "level": current_level,
            "title": current_title,
            "content": "\n".join(current_content),
        })
    return ParsedDocument(title=filename, sections=sections)


def _parse_html(content: bytes, filename: str) -> ParsedDocument:
    from bs4 import BeautifulSoup
    text = content.decode("utf-8", errors="ignore")
    soup = BeautifulSoup(text, "html.parser")
    # Remove script and style elements
    for script in soup(["script", "style"]):
        script.decompose()
    sections = []
    current_level = 0
    current_title = filename
    current_content = []
    for elem in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p"]):
        if elem.name.startswith("h"):
            if current_content or not sections:
                sections.append({
                    "level": current_level,
                    "title": current_title,
                    "content": "\n".join(current_content),
                })
            try:
                current_level = int(elem.name[1:])
            except ValueError:
                current_level = 0
            current_title = elem.get_text(strip=True)
            current_content = []
        else:
            txt = elem.get_text(strip=True)
            if txt:
                current_content.append(txt)
    if current_content or not sections:
        sections.append({
            "level": current_level,
            "title": current_title,
            "content": "\n".join(current_content),
        })
    return ParsedDocument(title=filename, sections=sections)


def _parse_csv(content: bytes, filename: str) -> ParsedDocument:
    text = content.decode("utf-8")
    reader = csv.reader(io.StringIO(text))
    rows = [" | ".join(row) for row in reader]
    return ParsedDocument(
        title=filename,
        sections=[{"level": 0, "title": filename, "content": "\n".join(rows)}],
    )


_MIME_HANDLERS = {
    "text/plain": _parse_txt,
    "application/pdf": _parse_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": _parse_docx,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": _parse_xlsx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": _parse_pptx,
    "text/markdown": _parse_md,
    "text/html": _parse_html,
    "text/csv": _parse_csv,
}


def parse_document(content: bytes, mime_type: str, filename: str) -> ParsedDocument:
    handler = _MIME_HANDLERS.get(mime_type)
    if handler is None:
        raise DocumentParseException(f"Unsupported MIME type: {mime_type}")
    try:
        return handler(content, filename)
    except Exception as exc:
        logger.exception("Document parsing failed for %s", filename)
        raise DocumentParseException(f"Failed to parse {filename}: {exc}") from exc
