import io
import pytest
from app.services.document_parser import parse_document, ParsedDocument, DocumentParseException


def test_parse_txt():
    content = b"Hello world\n\nThis is a test."
    result = parse_document(content, "text/plain", "test.txt")
    assert isinstance(result, ParsedDocument)
    assert result.title == "test.txt"
    assert len(result.sections) == 1
    assert result.sections[0]["level"] == 0
    assert "Hello world" in result.sections[0]["content"]


def test_parse_unsupported_mime_type():
    with pytest.raises(DocumentParseException, match="Unsupported MIME type"):
        parse_document(b"data", "application/unknown", "test.bin")


def test_parse_pdf():
    from pypdf import PdfWriter
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    # pypdf doesn't easily add text to blank pages; we test the extraction path at minimum
    buf = io.BytesIO()
    writer.write(buf)
    result = parse_document(buf.getvalue(), "application/pdf", "test.pdf")
    assert result.title == "test.pdf"
    assert len(result.sections) >= 1


def test_parse_docx():
    from docx import Document
    doc = Document()
    doc.add_heading("Heading 1", level=1)
    doc.add_paragraph("Paragraph text")
    buf = io.BytesIO()
    doc.save(buf)
    result = parse_document(buf.getvalue(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "test.docx")
    assert result.title == "test.docx"
    assert len(result.sections) >= 1
    assert any("Paragraph text" in s["content"] for s in result.sections)


def test_parse_markdown():
    content = b"# Title\n\nSome content.\n\n## Subtitle\n\nMore content."
    result = parse_document(content, "text/markdown", "test.md")
    assert result.title == "test.md"
    assert len(result.sections) >= 2
    assert any(s["title"] == "Title" for s in result.sections)
    assert any(s["title"] == "Subtitle" for s in result.sections)


def test_parse_html():
    content = b"<html><body><h1>Title</h1><p>Content</p></body></html>"
    result = parse_document(content, "text/html", "test.html")
    assert result.title == "test.html"
    assert len(result.sections) >= 1
    assert any("Content" in s["content"] for s in result.sections)


def test_parse_csv():
    content = b"name,age\nAlice,30\nBob,25"
    result = parse_document(content, "text/csv", "test.csv")
    assert result.title == "test.csv"
    assert len(result.sections) == 1
    assert "Alice" in result.sections[0]["content"]


def test_parse_xlsx():
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Header"
    ws["A2"] = "Value"
    buf = io.BytesIO()
    wb.save(buf)
    result = parse_document(buf.getvalue(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "test.xlsx")
    assert result.title == "test.xlsx"
    assert len(result.sections) >= 1
    assert any("Value" in s["content"] for s in result.sections)


def test_parse_pptx():
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Slide Title"
    buf = io.BytesIO()
    prs.save(buf)
    result = parse_document(buf.getvalue(), "application/vnd.openxmlformats-officedocument.presentationml.presentation", "test.pptx")
    assert result.title == "test.pptx"
    assert len(result.sections) >= 1
    assert any("Slide Title" in s["content"] for s in result.sections)
